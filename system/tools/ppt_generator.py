#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Copyright (c) 2026 BACH Contributors

Permission is hereby granted, free of charge, to any person obtaining a copy
of this software and associated documentation files (the "Software"), to deal
in the Software without restriction, including without limitation the rights
to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
copies of the Software, and to permit persons to whom the Software is
furnished to do so, subject to the following conditions:

The above copyright notice and this permission notice shall be included in all
copies or substantial portions of the Software.

THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
SOFTWARE.
"""

"""
BACH Zugangs-Übersicht PowerPoint Generator
Erstellt visuelle Präsentation der drei BACH-Zugänge
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import sys

def create_title_slide(prs):
    """Titelfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "BACH"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Untertitel
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Persönliches KI-Assistenz-System\nDie drei Zugänge"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_overview_slide(prs):
    """Übersichtsfolie - Vier Zugänge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Die vier Zugänge zu BACH"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Box 1: CLI
    cli_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(2.2), Inches(2.2))
    cli_box.fill.solid()
    cli_box.fill.fore_color.rgb = RGBColor(45, 90, 135)
    cli_box.line.color.rgb = RGBColor(255, 255, 255)
    cli_box.line.width = Pt(2)

    cli_text = cli_box.text_frame
    cli_text.text = "1️⃣ CLI\n\n⚡ Schnell\n🔄 Automation\n🔌 Hooks"
    for paragraph in cli_text.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Box 2: GUI
    gui_box = slide.shapes.add_shape(1, Inches(3), Inches(1.2), Inches(2.2), Inches(2.2))
    gui_box.fill.solid()
    gui_box.fill.fore_color.rgb = RGBColor(90, 135, 45)
    gui_box.line.color.rgb = RGBColor(255, 255, 255)
    gui_box.line.width = Pt(2)

    gui_text = gui_box.text_frame
    gui_text.text = "2️⃣ GUI\n\n📊 Visual\n🔴 Live\n🔌 Hooks"
    for paragraph in gui_text.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Box 3: Bridge
    bridge_box = slide.shapes.add_shape(1, Inches(5.5), Inches(1.2), Inches(2.2), Inches(2.2))
    bridge_box.fill.solid()
    bridge_box.fill.fore_color.rgb = RGBColor(135, 45, 90)
    bridge_box.line.color.rgb = RGBColor(255, 255, 255)
    bridge_box.line.width = Pt(2)

    bridge_text = bridge_box.text_frame
    bridge_text.text = "3️⃣ BRIDGE\n\n💾 24h\n🔗 MCP\n🔌 Hooks"
    for paragraph in bridge_text.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Box 4: Guest (NEU)
    guest_box = slide.shapes.add_shape(1, Inches(8), Inches(1.2), Inches(2.2), Inches(2.2))
    guest_box.fill.solid()
    guest_box.fill.fore_color.rgb = RGBColor(212, 160, 23)
    guest_box.line.color.rgb = RGBColor(255, 255, 255)
    guest_box.line.width = Pt(2)

    guest_text = guest_box.text_frame
    guest_text.text = "4️⃣ GUEST\n\n📖 Read-Only\n🌍 Portabel\n❌ Keine Hooks"
    for paragraph in guest_text.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Hook-System (unten)
    hook_box = slide.shapes.add_shape(1, Inches(0.5), Inches(3.7), Inches(7.2), Inches(1.5))
    hook_box.fill.solid()
    hook_box.fill.fore_color.rgb = RGBColor(85, 85, 85)
    hook_box.line.color.rgb = RGBColor(255, 255, 255)
    hook_box.line.width = Pt(2)

    hook_text = hook_box.text_frame
    hook_text.text = "🔌 HOOK-SYSTEM (Feature-Layer)\nVerfügbar in: CLI, GUI, Bridge | Nicht in: Guest"
    for paragraph in hook_text.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_user_types_slide(prs):
    """Nutzertypen & Empfehlungen"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Wer nutzt welchen Zugang?"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Power User → CLI
    user1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.2), Inches(1))
    user1_frame = user1_box.text_frame
    user1_frame.text = "🔧 Power User"
    for p in user1_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    rec1_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(2.2), Inches(0.7))
    rec1_frame = rec1_box.text_frame
    rec1_frame.text = "→ CLI"
    for p in rec1_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(45, 90, 135)
        p.alignment = PP_ALIGN.CENTER

    # Standard User → GUI
    user2_box = slide.shapes.add_textbox(Inches(3), Inches(1.2), Inches(2.2), Inches(1))
    user2_frame = user2_box.text_frame
    user2_frame.text = "👤 Standard User"
    for p in user2_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    rec2_box = slide.shapes.add_textbox(Inches(3), Inches(2.3), Inches(2.2), Inches(0.7))
    rec2_frame = rec2_box.text_frame
    rec2_frame.text = "→ GUI"
    for p in rec2_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(90, 135, 45)
        p.alignment = PP_ALIGN.CENTER

    # Claude User → Bridge
    user3_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(2.2), Inches(1))
    user3_frame = user3_box.text_frame
    user3_frame.text = "🤖 Claude User"
    for p in user3_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    rec3_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(2.2), Inches(0.7))
    rec3_frame = rec3_box.text_frame
    rec3_frame.text = "→ BRIDGE"
    for p in rec3_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(135, 45, 90)
        p.alignment = PP_ALIGN.CENTER

    # Guest User → Guest (NEU)
    user4_box = slide.shapes.add_textbox(Inches(8), Inches(1.2), Inches(2.2), Inches(1))
    user4_frame = user4_box.text_frame
    user4_frame.text = "🎓 Externe/Gäste"
    for p in user4_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    rec4_box = slide.shapes.add_textbox(Inches(8), Inches(2.3), Inches(2.2), Inches(0.7))
    rec4_frame = rec4_box.text_frame
    rec4_frame.text = "→ GUEST"
    for p in rec4_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(212, 160, 23)
        p.alignment = PP_ALIGN.CENTER

    # Features unten
    features_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    features_frame = features_box.text_frame
    features_frame.text = """CLI: Schnell, Automation, Hooks, Alle Handler
GUI: Visualisierung, Live-Updates, Monitoring
BRIDGE: 24h-Session, MCP, Bidirektional, Skills
GUEST: Keine Installation, Portabel, Read-Only, Kein System-Zugriff"""
    for p in features_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.space_after = Pt(8)

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_backend_slide(prs):
    """Gemeinsames Backend"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Gemeinsames Backend - Die Infrastruktur"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Drei Zugänge oben
    cli_small = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.5), Inches(0.5))
    cli_small.text_frame.text = "CLI"
    cli_small.text_frame.paragraphs[0].font.size = Pt(20)
    cli_small.text_frame.paragraphs[0].font.color.rgb = RGBColor(45, 90, 135)
    cli_small.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    gui_small = slide.shapes.add_textbox(Inches(3.75), Inches(1.2), Inches(2.5), Inches(0.5))
    gui_small.text_frame.text = "GUI"
    gui_small.text_frame.paragraphs[0].font.size = Pt(20)
    gui_small.text_frame.paragraphs[0].font.color.rgb = RGBColor(90, 135, 45)
    gui_small.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    bridge_small = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(2.5), Inches(0.5))
    bridge_small.text_frame.text = "BRIDGE"
    bridge_small.text_frame.paragraphs[0].font.size = Pt(20)
    bridge_small.text_frame.paragraphs[0].font.color.rgb = RGBColor(135, 45, 90)
    bridge_small.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Backend Box
    backend_box = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(8), Inches(3))
    backend_box.fill.solid()
    backend_box.fill.fore_color.rgb = RGBColor(68, 68, 68)
    backend_box.line.color.rgb = RGBColor(255, 255, 255)
    backend_box.line.width = Pt(3)

    backend_text = backend_box.text_frame
    backend_text.text = "🏗️ GEMEINSAMES BACKEND\n\n🗄️ bach.db (SQLite3, 32+ Tables)\n⚙️ Handler-System (Auto-Discovery)\n📚 Skill-System (32+ Skills + Agents)\n📝 Auto-Logger (Extended Logs)"
    for p in backend_text.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_comparison_slide(prs):
    """Vergleichsmatrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Vergleich der Zugänge"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Tabelle als Text (python-pptx hat eingeschränkten Table-Support)
    table_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(9.4), Inches(4.5))
    table_frame = table_box.text_frame
    table_frame.text = """KRITERIUM        CLI           GUI           BRIDGE        GUEST
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🎯 Zielgruppe    Power User    Alle          Claude User   Externe
💻 Interface     Terminal      Browser       MCP           Skill-Upload
💾 Zustand       Stateless     Session       Persistent    Stateless
🔗 System        ✅ Voll       ✅ Voll       ✅ Voll       ❌ Kein
⚡ Performance   ⭐⭐⭐⭐⭐     ⭐⭐⭐         ⭐⭐⭐⭐       ⭐⭐⭐⭐⭐
🔄 Automation    ⭐⭐⭐⭐⭐     ⭐⭐           ⭐⭐⭐         ❌
📊 Visual        ⭐             ⭐⭐⭐⭐⭐     ⭐⭐           ⭐
📱 Portabel      ⚠️ SSH        ⚠️ Local     ❌ Local      ✅ Überall
🔌 Hooks         ✅             ✅             ✅             ❌
🛠️ Install       ✅ Nötig      ✅ Nötig      ✅ Nötig      ❌ Nicht"""

    for p in table_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_scenarios_slide(prs):
    """Empfohlene Szenarien"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Wann welchen Zugang nutzen?"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Szenarien
    scenarios_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    scenarios_frame = scenarios_box.text_frame
    scenarios_frame.text = """📝 Schnell Task erstellen                 → CLI
📊 System-Übersicht anzeigen              → GUI
🤖 Mit Claude Code arbeiten               → BRIDGE
🔄 Automatisierung/Scripting              → CLI
📈 Monitoring & Dashboards                → GUI
🎓 Interaktive Agenten-Arbeit             → BRIDGE
🤝 Multi-Partner-Koordination             → CLI + Partner
🔍 Logs durchsuchen                       → GUI oder CLI
💡 BACH-Beratung ohne Installation        → GUEST
🌍 Unterwegs/ohne BACH-Zugriff            → GUEST
📖 Skill-Dokumentation nutzen             → GUEST"""

    for p in scenarios_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(10)

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_security_slide(prs):
    """Sicherheitskonzept"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Sicherheitsarchitektur"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Trust-Levels
    trust_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    trust_box.fill.solid()
    trust_box.fill.fore_color.rgb = RGBColor(68, 68, 68)
    trust_box.line.color.rgb = RGBColor(255, 255, 255)
    trust_box.line.width = Pt(2)

    trust_text = trust_box.text_frame
    trust_text.text = "🔐 Trust-Level System\n\n🥇 Goldstandard (100)\n✅ Trusted (80)\n⚠️ Untrusted (20)\n❌ Blacklist (0)"
    for p in trust_text.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Filesystem Protection
    fs_box = slide.shapes.add_shape(1, Inches(5.5), Inches(1.2), Inches(4), Inches(2))
    fs_box.fill.solid()
    fs_box.fill.fore_color.rgb = RGBColor(68, 68, 68)
    fs_box.line.color.rgb = RGBColor(255, 255, 255)
    fs_box.line.width = Pt(2)

    fs_text = fs_box.text_frame
    fs_text.text = "📁 Filesystem Protection\n\n📁 Core → Schreibgeschützt\n📄 Templates → Read-Only\n✏️ User → Vollzugriff\n💾 Auto-Backup"
    for p in fs_text.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Sandboxing
    sandbox_box = slide.shapes.add_shape(1, Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    sandbox_box.fill.solid()
    sandbox_box.fill.fore_color.rgb = RGBColor(68, 68, 68)
    sandbox_box.line.color.rgb = RGBColor(255, 255, 255)
    sandbox_box.line.width = Pt(2)

    sandbox_text = sandbox_box.text_frame
    sandbox_text.text = "🛡️ Sandboxing-Regeln\n\n❌ bridge_daemon.py → NICHT editierbar\n❌ hub/_services/ → NICHT änderbar\n✅ Skills → Lesbar, exportierbar\n✅ Tools → Mit Whitelist-Check"
    for p in sandbox_text.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def create_final_slide(prs):
    """Abschlussfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Zusammenfassung
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    summary_frame = summary_box.text_frame
    summary_frame.text = "BACH - Ein System, Vier Wege\n\n✅ Flexibilität durch vier Zugänge\n✅ Konsistenz durch gemeinsames Backend\n✅ Portabilität durch Guest-Modus\n✅ Erweiterbarkeit durch Hook-System\n✅ Parallelität - Alle Zugänge nutzbar\n\n🚪 bach.bat - Ihr zentraler Einstieg"

    for p in summary_frame.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(16)

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)

def main():
    """Hauptfunktion"""
    # Präsentation erstellen
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Folien erstellen
    print("Erstelle Titelfolie...")
    create_title_slide(prs)

    print("Erstelle Übersichtsfolie...")
    create_overview_slide(prs)

    print("Erstelle Nutzertypen-Folie...")
    create_user_types_slide(prs)

    print("Erstelle Backend-Folie...")
    create_backend_slide(prs)

    print("Erstelle Vergleichsfolie...")
    create_comparison_slide(prs)

    print("Erstelle Szenarien-Folie...")
    create_scenarios_slide(prs)

    print("Erstelle Security-Folie...")
    create_security_slide(prs)

    print("Erstelle Abschlussfolie...")
    create_final_slide(prs)

    # Speichern
    output_path = "docs/BACH_Zugangs_Uebersicht.pptx"
    prs.save(output_path)
    print(f"\n✅ Präsentation erstellt: {output_path}")
    print(f"📊 {len(prs.slides)} Folien")

if __name__ == "__main__":
    main()
